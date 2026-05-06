from pathlib import Path
from textwrap import dedent

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DocxPt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "defense_output"
OUT.mkdir(exist_ok=True)

TITLE = "基于 MindSpore 的上下文 GAN 图像修复与超分辨率联合处理系统设计"
SUBTITLE = "CRA 图像修复 × SRGAN 超分辨率 × Tkinter 可视化交互"
AUTHOR = "苗玉鑫｜智能科学与工程学院｜网络工程"


COLORS = {
    "navy": RGBColor(15, 23, 42),
    "slate": RGBColor(51, 65, 85),
    "muted": RGBColor(100, 116, 139),
    "blue": RGBColor(37, 99, 235),
    "cyan": RGBColor(6, 182, 212),
    "green": RGBColor(22, 163, 74),
    "orange": RGBColor(245, 158, 11),
    "red": RGBColor(220, 38, 38),
    "light": RGBColor(248, 250, 252),
    "line": RGBColor(226, 232, 240),
    "white": RGBColor(255, 255, 255),
}


def font(size=20, bold=False):
    return {"name": "Microsoft YaHei", "size": Pt(size), "bold": bold}


def set_run(run, size=20, color="slate", bold=False):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]


def add_textbox(slide, text, x, y, w, h, size=20, color="slate", bold=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, color, bold)
    return box


def add_title(slide, title, eyebrow=None):
    if eyebrow:
        add_textbox(slide, eyebrow, 0.65, 0.34, 5.5, 0.28, 10, "blue", True)
    add_textbox(slide, title, 0.65, 0.58, 8.2, 0.55, 24, "navy", True)
    line = slide.shapes.add_shape(1, Inches(0.65), Inches(1.22), Inches(1.0), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["blue"]
    line.line.color.rgb = COLORS["blue"]


def add_footer(slide, idx):
    add_textbox(slide, f"{idx:02d}", 12.2, 7.05, 0.45, 0.25, 9, "muted", True, PP_ALIGN.RIGHT)
    add_textbox(slide, "CRA-SRGAN 联合图像复原系统答辩", 0.65, 7.05, 4.0, 0.25, 9, "muted")


def bullet_list(slide, items, x, y, w, h, size=15, gap=0.33):
    top = y
    for item in items:
        box = slide.shapes.add_textbox(Inches(x), Inches(top), Inches(w), Inches(gap))
        tf = box.text_frame
        tf.clear()
        tf.margin_left = Inches(0.02)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"• {item}"
        set_run(run, size, "slate")
        top += gap


def add_card(slide, title, body, x, y, w, h, accent="blue"):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[accent]
    bar.line.color.rgb = COLORS[accent]
    add_textbox(slide, title, x + 0.2, y + 0.15, w - 0.35, 0.3, 15, "navy", True)
    add_textbox(slide, body, x + 0.2, y + 0.55, w - 0.35, h - 0.68, 12.5, "slate")


def add_image(slide, path, x, y, w, h):
    path = Path(path)
    if not path.exists():
        add_card(slide, "缺少图片", str(path), x, y, w, h, "red")
        return
    with Image.open(path) as img:
        iw, ih = img.size
    ratio = min(w / iw, h / ih)
    real_w = iw * ratio
    real_h = ih * ratio
    left = x + (w - real_w) / 2
    top = y + (h - real_h) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(real_w), Inches(real_h))


def add_metric(slide, value, label, x, y, color="blue"):
    add_textbox(slide, value, x, y, 1.9, 0.42, 22, color, True, PP_ALIGN.CENTER)
    add_textbox(slide, label, x, y + 0.46, 1.9, 0.32, 10.5, "muted", False, PP_ALIGN.CENTER)


def make_comparison_image():
    paths = [
        ROOT / "output_demo_final" / "demo_assets" / "demo_input.png",
        ROOT / "output_demo_final" / "demo_assets" / "demo_mask.png",
        ROOT / "output_demo_final" / "demo_input_repaired.png",
        ROOT / "output_demo_final" / "demo_input_enhanced.png",
    ]
    labels = ["输入图像", "修复掩码", "CRA 修复结果", "SRGAN 增强结果"]
    imgs = []
    for p in paths:
        if p.exists():
            imgs.append(Image.open(p).convert("RGB"))
        else:
            imgs.append(Image.new("RGB", (256, 192), (240, 244, 248)))
    tile_w, tile_h = 280, 210
    canvas = Image.new("RGB", (tile_w * 4, tile_h + 48), (248, 250, 252))
    draw = ImageDraw.Draw(canvas)
    try:
        fnt = ImageFont.truetype("msyh.ttc", 22)
    except Exception:
        fnt = ImageFont.load_default()
    for i, img in enumerate(imgs):
        img.thumbnail((tile_w - 24, tile_h - 45), Image.LANCZOS)
        x = i * tile_w + (tile_w - img.width) // 2
        y = 14
        canvas.paste(img, (x, y))
        draw.text((i * tile_w + 22, tile_h - 24), labels[i], fill=(51, 65, 85), font=fnt)
        if i < 3:
            draw.text((i * tile_w + tile_w - 12, tile_h // 2), "→", fill=(37, 99, 235), font=fnt)
    out = OUT / "demo_comparison.png"
    canvas.save(out)
    return out


def make_gui_mock():
    out = OUT / "gui_mock.png"
    W, H = 1200, 760
    img = Image.new("RGB", (W, H), (248, 250, 252))
    d = ImageDraw.Draw(img)
    try:
        title = ImageFont.truetype("msyh.ttc", 34)
        head = ImageFont.truetype("msyh.ttc", 22)
        normal = ImageFont.truetype("msyh.ttc", 18)
        small = ImageFont.truetype("msyh.ttc", 15)
    except Exception:
        title = head = normal = small = ImageFont.load_default()
    d.rounded_rectangle((32, 28, W - 32, H - 28), radius=18, fill=(255, 255, 255), outline=(226, 232, 240), width=2)
    d.text((62, 56), "图像修复与超分辨率处理工具", fill=(15, 23, 42), font=title)
    d.rounded_rectangle((910, 58, 1110, 98), radius=8, fill=(219, 234, 254))
    d.text((938, 67), "classic 后端就绪", fill=(37, 99, 235), font=normal)
    sections = [
        ("文件选择", 64, 126, 1090, 110),
        ("模型配置", 64, 254, 1090, 96),
        ("参数配置", 64, 368, 1090, 82),
        ("图像预览", 64, 468, 710, 190),
        ("运行日志", 792, 468, 362, 190),
    ]
    for name, x, y, w, h in sections:
        d.rounded_rectangle((x, y, x + w, y + h), radius=10, fill=(248, 250, 252), outline=(226, 232, 240), width=2)
        d.text((x + 18, y + 12), name, fill=(15, 23, 42), font=head)
    for i, label in enumerate(["输入图像", "掩码图像", "输出目录"]):
        y = 166 + i * 30
        d.text((90, y), label, fill=(71, 85, 105), font=small)
        d.rounded_rectangle((190, y - 4, 850, y + 22), radius=5, fill=(255, 255, 255), outline=(203, 213, 225))
        d.rounded_rectangle((875, y - 5, 940, y + 23), radius=5, fill=(37, 99, 235))
        d.text((890, y), "浏览", fill=(255, 255, 255), font=small)
    for i, label in enumerate(["CRA 模型", "SRGAN 模型"]):
        y = 294 + i * 30
        d.text((90, y), label, fill=(71, 85, 105), font=small)
        d.rounded_rectangle((190, y - 4, 850, y + 22), radius=5, fill=(255, 255, 255), outline=(203, 213, 225))
    for i, label in enumerate(["生成演示输入", "检查掩码", "加载模型", "开始处理"]):
        x = 90 + i * 160
        d.rounded_rectangle((x, 675, x + 130, 720), radius=8, fill=(37, 99, 235) if i == 3 else (241, 245, 249),
                            outline=(203, 213, 225))
        d.text((x + 18, 687), label, fill=(255, 255, 255) if i == 3 else (51, 65, 85), font=small)
    d.text((90, 408), "后端: auto    设备: CPU    放大倍数: 2×    输入尺寸: 512", fill=(51, 65, 85), font=normal)
    for i, lab in enumerate(["输入图", "掩码", "结果"]):
        x = 100 + i * 220
        d.rounded_rectangle((x, 520, x + 170, 625), radius=8, fill=(226, 232, 240))
        d.text((x + 55, 555), lab, fill=(100, 116, 139), font=head)
    for i, line in enumerate(["22:51:17 - 加载图像成功", "22:51:17 - OpenCV 修复完成", "22:51:17 - 超分增强完成"]):
        d.text((820, 520 + i * 36), line, fill=(51, 65, 85), font=small)
    img.save(out)
    return out


def build_deck():
    comparison = make_comparison_image()
    gui_mock = make_gui_mock()
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = []
    for _ in range(16):
        slides.append(prs.slides.add_slide(blank))

    # 1
    s = slides[0]
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["navy"]
    add_textbox(s, TITLE, 0.72, 1.45, 11.6, 1.45, 32, "white", True)
    add_textbox(s, SUBTITLE, 0.78, 3.05, 9.8, 0.45, 18, "cyan")
    add_textbox(s, AUTHOR + "｜2026 年本科毕业设计答辩", 0.8, 6.35, 8.5, 0.32, 13, "light")
    for i, color in enumerate(["blue", "cyan", "green"]):
        shp = s.shapes.add_shape(1, Inches(10.2 + i * 0.38), Inches(5.65 - i * 0.28), Inches(1.5), Inches(0.16))
        shp.fill.solid()
        shp.fill.fore_color.rgb = COLORS[color]
        shp.line.color.rgb = COLORS[color]

    # 2
    s = slides[1]
    add_title(s, "答辩内容安排", "CONTENTS")
    agenda = [
        ("01", "研究背景与问题定义", "为什么要把图像修复和超分辨率联合起来"),
        ("02", "核心模型与系统设计", "CRA 负责补全内容，SRGAN 负责增强细节"),
        ("03", "工程实现与交互界面", "模块化、分块融合、GUI 与批处理"),
        ("04", "测试结果与总结展望", "质量、速度、可用性与后续优化方向"),
    ]
    for i, (num, title, body) in enumerate(agenda):
        add_card(s, f"{num}  {title}", body, 0.9 + (i % 2) * 5.7, 1.7 + (i // 2) * 1.75, 5.1, 1.2, ["blue", "cyan", "green", "orange"][i])
    add_footer(s, 2)

    # 3
    s = slides[2]
    add_title(s, "研究背景：真实图像常常是复合退化", "BACKGROUND")
    bullet_list(s, [
        "真实场景中图像会同时出现遮挡、破损、模糊、低分辨率等问题",
        "单一图像修复方法关注“补内容”，但无法提升整体清晰度",
        "单一超分辨率方法关注“提分辨率”，但会放大破损区域和伪影",
        "因此需要一条连续流程：先恢复结构，再增强纹理"
    ], 0.85, 1.55, 5.6, 2.1, 16, 0.42)
    add_card(s, "一句话概括", "本课题要解决的不是单张图片变清楚，而是让“有缺损、又不清楚”的图片完成联合复原。", 0.85, 4.2, 5.6, 1.25, "blue")
    add_image(s, ROOT / "output_demo_final" / "demo_input_enhanced.png", 7.25, 1.35, 4.7, 4.0)
    add_footer(s, 3)

    # 4
    s = slides[3]
    add_title(s, "研究目标与本文工作", "OBJECTIVE")
    cards = [
        ("联合处理链路", "构建“输入图像 + 掩码 → CRA 修复 → SRGAN 增强 → 结果保存”的完整流程。", "blue"),
        ("统一工程封装", "将配置、图像读写、模型加载、推理、GUI 与批量处理拆成清晰模块。", "cyan"),
        ("可展示可运行", "提供命令行、GUI、测试脚本和可视化图表，便于论文展示和答辩演示。", "green"),
        ("稳定性优化", "加入输入校验、异常回退、日志提示、分块推理与重叠融合。", "orange"),
    ]
    for i, (t, b, c) in enumerate(cards):
        add_card(s, t, b, 0.85 + (i % 2) * 5.75, 1.55 + (i // 2) * 1.72, 5.25, 1.24, c)
    add_footer(s, 4)

    # 5
    s = slides[4]
    add_title(s, "总体架构：模型层、流程层、应用层分离", "ARCHITECTURE")
    add_image(s, ROOT / "模型流程图" / "项目总体架构图.png", 0.8, 1.35, 6.15, 4.75)
    bullet_list(s, [
        "模型层：CRA 图像修复模型、SRGAN 超分辨率生成器",
        "流程层：CombinedProcessor 统一调度修复与增强",
        "应用层：CLI、GUI、批处理脚本共享同一核心流程",
        "支撑层：配置管理、日志、检查点发现、图表生成"
    ], 7.35, 1.55, 4.9, 2.3, 15, 0.4)
    add_footer(s, 5)

    # 6
    s = slides[5]
    add_title(s, "核心处理流程：先修复，后增强", "PIPELINE")
    add_image(s, ROOT / "模型流程图" / "组合架构总览图.png", 0.85, 1.35, 5.5, 4.5)
    add_card(s, "为什么不直接超分？", "如果直接放大破损图像，缺损区域和边缘伪影也会被一起放大；先用 CRA 让内容结构完整，再用 SRGAN 补充高频细节，最终视觉效果更自然。", 6.95, 1.65, 5.35, 1.55, "blue")
    add_card(s, "统一输出", "系统保存两类结果：*_repaired.png 作为修复中间结果，*_enhanced.png 作为最终增强结果，方便对比和回溯。", 6.95, 3.65, 5.35, 1.3, "green")
    add_footer(s, 6)

    # 7
    s = slides[6]
    add_title(s, "CRA 图像修复：利用上下文补全缺损区域", "CRA INPAINTING")
    add_image(s, ROOT / "模型流程图" / "CRA模型架构图.png", 0.85, 1.35, 5.65, 4.75)
    bullet_list(s, [
        "输入：原始图像与二值掩码，白色区域表示需要修复",
        "粗修复阶段先恢复整体结构，避免结果语义断裂",
        "上下文注意力从未破损区域寻找可参考的纹理和结构",
        "残差聚合把高频细节迁移回缺损区域，提升连贯性"
    ], 7.0, 1.55, 5.1, 2.2, 15, 0.42)
    add_footer(s, 7)

    # 8
    s = slides[7]
    add_title(s, "SRGAN 超分辨率：提升分辨率与视觉细节", "SRGAN SUPER-RESOLUTION")
    add_image(s, ROOT / "模型流程图" / "SRGAN模型架构图.png", 0.75, 1.35, 5.8, 4.7)
    bullet_list(s, [
        "生成器由卷积层、残差块和子像素上采样模块组成",
        "残差学习保留主体结构，降低深层网络训练难度",
        "子像素卷积负责把低分辨率特征重建为高分辨率图像",
        "对抗训练更关注感知质量，适合恢复纹理细节"
    ], 7.0, 1.55, 5.1, 2.2, 15, 0.42)
    add_footer(s, 8)

    # 9
    s = slides[8]
    add_title(s, "工程重点：大图分块推理与重叠融合", "ENGINEERING")
    add_image(s, ROOT / "模型流程图" / "CRA和SRGAN联合推理流程图.png", 0.75, 1.4, 6.0, 4.8)
    add_metric(s, "14.2GB", "整体处理显存", 7.1, 1.7, "red")
    add_metric(s, "5.3GB", "分块处理显存", 9.3, 1.7, "green")
    add_metric(s, "≈63%", "显存占用下降", 11.1, 1.7, "blue")
    add_card(s, "重叠融合策略", "按块裁剪 → 单块增强 → 构建权重图 → 重叠区域渐变加权 → 归一化输出。这样可以降低显存压力，同时减轻块边界拼接痕迹。", 7.15, 3.05, 5.1, 1.7, "cyan")
    add_footer(s, 9)

    # 10
    s = slides[9]
    add_title(s, "系统实现：模块化代码组织", "IMPLEMENTATION")
    modules = [
        ("ModelConfig", "统一维护后端、设备、倍率、检查点与 .env/JSON 配置"),
        ("ImageProcessor", "图像读取、掩码二值化、尺寸匹配、结果保存"),
        ("CRAModel / SRGANModel", "封装深度模型加载、推理和后处理逻辑"),
        ("Classic fallback", "无 MindSpore 时使用 OpenCV 修复与插值增强保证可运行"),
        ("CombinedProcessor", "统一调度修复、增强、异常回退与输出命名"),
        ("ImageRepairGUI", "文件选择、参数设置、日志反馈、结果预览与演示输入"),
    ]
    for i, (t, b) in enumerate(modules):
        add_card(s, t, b, 0.85 + (i % 2) * 5.85, 1.38 + (i // 2) * 1.35, 5.35, 0.94, ["blue", "cyan", "green", "orange", "blue", "cyan"][i])
    add_footer(s, 10)

    # 11
    s = slides[10]
    add_title(s, "可视化交互界面：降低使用门槛", "GUI")
    add_image(s, gui_mock, 0.8, 1.35, 6.1, 4.9)
    bullet_list(s, [
        "文件选择：输入图像、掩码、输出目录",
        "模型配置：自动检测检查点，也支持 classic 后端直接演示",
        "参数设置：后端、设备、超分倍数、输入尺寸",
        "实时反馈：运行日志、进度条、结果预览",
        "辅助能力：生成演示输入、检查掩码有效性"
    ], 7.25, 1.55, 4.9, 2.7, 14.5, 0.38)
    add_footer(s, 11)

    # 12
    s = slides[11]
    add_title(s, "端到端处理效果展示", "DEMO RESULT")
    add_image(s, comparison, 0.8, 1.35, 11.9, 3.3)
    add_card(s, "展示逻辑", "输入图像存在缺损区域，掩码标出待修复位置；系统先输出修复中间图，再输出分辨率增强后的最终结果。答辩时可现场运行 GUI 的“生成演示输入 → 加载模型 → 开始处理”。", 1.2, 5.15, 10.9, 1.05, "green")
    add_footer(s, 12)

    # 13
    s = slides[12]
    add_title(s, "测试结果：质量、感知和速度指标", "EVALUATION")
    add_image(s, ROOT / "visual_results" / "PSNR性能对比图.png", 0.75, 1.38, 3.95, 3.7)
    add_image(s, ROOT / "visual_results" / "超分辨率性能对比图.png", 4.7, 1.38, 3.95, 3.7)
    add_image(s, ROOT / "visual_results" / "实时处理性能图.png", 8.65, 1.38, 3.95, 3.7)
    add_metric(s, "29.45 dB", "BSD100 PSNR", 1.1, 5.7, "blue")
    add_metric(s, "0.8623", "SSIM", 3.35, 5.7, "green")
    add_metric(s, "3.92", "FID", 5.8, 5.7, "orange")
    add_metric(s, "0.75s", "8K 平均处理", 8.25, 5.7, "cyan")
    add_metric(s, "5.3GB", "8K 分块显存", 10.55, 5.7, "green")
    add_footer(s, 13)

    # 14
    s = slides[13]
    add_title(s, "系统测试与稳定性验证", "TESTING")
    add_card(s, "功能测试", "覆盖输入读取、掩码预处理、模型加载、图像修复、超分增强、结果保存和图表生成。", 0.9, 1.55, 5.35, 1.2, "blue")
    add_card(s, "异常处理", "模型路径错误、掩码缺失、图像格式不支持、增强阶段异常时，系统通过日志和弹窗给出提示。", 0.9, 3.0, 5.35, 1.2, "orange")
    add_card(s, "交互验证", "GUI 使用线程化处理，模型加载和图像处理期间界面保持响应，并实时更新状态。", 6.95, 1.55, 5.35, 1.2, "cyan")
    add_card(s, "批量场景", "批处理脚本支持目录遍历、掩码匹配、失败样本统计和平均耗时输出。", 6.95, 3.0, 5.35, 1.2, "green")
    add_card(s, "本地可运行性补强", "在当前 Python 3.13 环境 MindSpore 不可安装时，增加 classic fallback，确保答辩演示可跑通。", 3.1, 4.62, 7.2, 1.05, "blue")
    add_footer(s, 14)

    # 15
    s = slides[14]
    add_title(s, "创新点、不足与后续工作", "SUMMARY")
    add_card(s, "主要创新与完成度", "将 CRA 与 SRGAN 串联为完整工程链路；加入分块融合、自动检查点发现、GUI、批处理和掩码分析，形成可运行、可展示、可维护的系统。", 0.9, 1.5, 5.45, 1.6, "green")
    add_card(s, "当前不足", "串联式流程仍存在误差传递；复杂纹理和极端退化样本泛化验证不足；真实大规模测试集仍需扩充。", 6.9, 1.5, 5.45, 1.6, "orange")
    add_card(s, "后续方向", "探索更紧密的联合建模，引入轻量化注意力或扩散先验；增强 GUI 任务管理能力；扩展老照片、监控图像等真实应用数据集。", 0.9, 3.55, 11.45, 1.35, "blue")
    add_footer(s, 15)

    # 16
    s = slides[15]
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["navy"]
    add_textbox(s, "谢谢各位老师", 0.9, 2.15, 11.4, 0.75, 36, "white", True, PP_ALIGN.CENTER)
    add_textbox(s, "请批评指正", 0.9, 3.05, 11.4, 0.45, 20, "cyan", False, PP_ALIGN.CENTER)
    add_textbox(s, "Q & A", 0.9, 4.15, 11.4, 0.65, 30, "light", True, PP_ALIGN.CENTER)

    deck_path = OUT / "基于MindSpore的上下文GAN图像修复与超分辨率联合处理系统_答辩PPT.pptx"
    prs.save(deck_path)
    return deck_path


SCRIPT = dedent(
    """
    # 答辩演讲稿（约 15 分钟）

    ## 第 1 页：标题页（约 40 秒）
    各位老师好，我是苗玉鑫。我的毕业设计题目是《基于 MindSpore 的上下文 GAN 图像修复与超分辨率联合处理系统设计》。本课题主要围绕图像复原中的两个常见问题展开：一个是图像内容存在遮挡或破损，另一个是图像分辨率不足、细节不清晰。我的工作是把 CRA 图像修复模型和 SRGAN 超分辨率模型进行工程化整合，形成一个从输入、修复、增强到结果展示的完整系统，并提供图形化界面，方便用户直接操作。

    ## 第 2 页：答辩内容安排（约 40 秒）
    本次答辩我会从四个部分展开。第一部分说明研究背景和问题定义，也就是为什么要做图像修复与超分辨率的联合处理。第二部分介绍核心模型和系统设计，重点说明 CRA 与 SRGAN 在系统中的分工。第三部分介绍工程实现，包括模块化封装、大图分块处理、GUI 交互和批量处理。最后介绍测试结果、系统不足和后续优化方向。

    ## 第 3 页：研究背景（约 1 分钟）
    在真实场景中，图像质量问题往往不是单一类型。例如安防监控图像可能既分辨率低，又存在遮挡；老照片可能既有划痕破损，又整体模糊；数字媒体素材也可能同时存在缺失区域和清晰度不足。传统方法如果只做图像修复，能够补全缺损区域，但不能提升整体分辨率；如果只做超分辨率，则会把破损区域和伪影一起放大。因此，本课题的核心出发点是：先恢复图像内容的完整性，再提升图像的清晰度和细节表现。

    ## 第 4 页：研究目标与本文工作（约 1 分钟）
    我的系统目标可以概括为四点。第一，构建一条完整的联合处理链路，从输入图像和掩码开始，先进行 CRA 修复，再进行 SRGAN 增强，最后保存结果。第二，对模型和流程进行统一封装，避免用户分别调用多个脚本。第三，系统不仅要能运行，还要能展示，所以实现了 GUI、批处理和结果可视化。第四，考虑工程稳定性，加入输入校验、异常处理、日志提示和分块推理等机制，使系统更接近一个完整可用的应用，而不只是模型实验代码。

    ## 第 5 页：总体架构（约 1 分钟）
    这一页展示的是系统总体架构。系统可以分为模型层、流程层、应用层和支撑层。模型层包括 CRA 图像修复模型和 SRGAN 超分辨率模型；流程层由 CombinedProcessor 统一调度，负责把修复结果交给增强模块；应用层包括命令行入口、GUI 图形界面和批处理脚本；支撑层包括配置管理、日志、检查点自动发现和图表生成。这样的分层设计使系统结构更清晰，后期如果替换模型或者增加新的后端，也不需要大范围改动应用层代码。

    ## 第 6 页：核心处理流程（约 1 分钟）
    系统采用“先修复、后增强”的两阶段流程。第一阶段，CRA 根据原图和掩码恢复缺损区域，让图像内容结构先变完整；第二阶段，SRGAN 对修复后的图像进行超分辨率增强，补充纹理和高频细节。这样做的好处是避免直接对破损图像进行放大，因为直接放大会使破损区域和边缘伪影更加明显。系统最终会保存两类结果：一个是修复中间结果，另一个是最终增强结果，便于用户进行对比分析。

    ## 第 7 页：CRA 图像修复模型（约 1 分 20 秒）
    CRA 是 Contextual Residual Aggregation，也就是上下文残差聚合模型。它在系统中负责解决“图像缺了什么内容”的问题。模型输入包括原始图像和二值掩码，白色区域表示需要修复的位置。CRA 的核心思想是利用未破损区域的上下文信息，对缺损区域进行结构和纹理补全。具体来说，模型先通过粗修复阶段恢复大致结构，再通过上下文注意力机制从周围区域寻找可参考的纹理信息，最后通过残差聚合将高频细节迁移回缺损区域。相比简单的邻域填充，这种方法更适合处理较大面积破损，因为它能够利用更远距离的上下文信息。

    ## 第 8 页：SRGAN 超分辨率模型（约 1 分 20 秒）
    SRGAN 是 Super-Resolution Generative Adversarial Network，也就是超分辨率生成对抗网络。它在系统中负责解决“图像不够清晰”的问题。SRGAN 的生成器由卷积层、多个残差块和子像素卷积上采样层组成。残差块有助于保留主体结构并稳定深层网络训练，子像素卷积负责把低分辨率特征重建为高分辨率图像。与普通插值方法相比，SRGAN 不只是让图片尺寸变大，而是通过对抗训练和感知质量约束，让结果更接近真实图像的纹理分布，因此更适合放在修复阶段之后，用来提升最终视觉质量。

    ## 第 9 页：大图分块推理与重叠融合（约 1 分 20 秒）
    在工程实现中，高分辨率图像会带来显存压力。如果把整张大图直接送入超分辨率网络，容易出现显存不足或运行失败。为了解决这个问题，我在增强模块中实现了分块推理与重叠融合机制。系统会根据图像尺寸和块大小自动切分图像，每个子块单独推理；在输出阶段，对重叠区域使用渐变权重进行融合，最后再归一化生成完整结果。这样既降低了显存占用，又能减少分块边界处的拼接痕迹。论文中的数据表明，8K 场景下显存可由 14.2GB 降到 5.3GB，降低约 63%。

    ## 第 10 页：系统实现与模块化代码（约 1 分钟）
    在代码实现上，我重点做了模块化封装。ModelConfig 负责统一维护后端、设备、超分倍率和检查点路径；ImageProcessor 负责图像读取、掩码二值化和结果保存；CRAModel 与 SRGANModel 分别封装深度模型加载和推理逻辑；CombinedProcessor 负责联合调度。同时，为了保证当前环境下也可以顺利演示，我增加了 classic fallback 后端。当 MindSpore 不可用时，系统会使用 OpenCV inpaint 和插值增强完成端到端流程，保证答辩和展示时系统不会因为环境问题完全不可运行。

    ## 第 11 页：GUI 可视化交互界面（约 1 分钟）
    图形界面基于 Tkinter 实现，主要目的是降低使用门槛。用户可以在界面中选择输入图像、掩码图像和输出目录，也可以配置模型路径、处理后端、运行设备、超分倍数和输入尺寸。界面还增加了运行日志、进度条和结果预览功能，用户可以看到输入图、掩码图和处理结果。为了便于答辩展示，我还增加了“生成演示输入”和“检查掩码”功能，点击后即可自动生成示例素材并检查修复区域比例。

    ## 第 12 页：端到端处理效果（约 50 秒）
    这一页展示的是端到端处理效果。左侧是输入图像，中间是掩码，表示需要修复的位置，然后系统输出 CRA 修复结果，最后输出 SRGAN 增强结果。这个流程体现了系统的主线：先让内容完整，再让图像更清晰。答辩现场如果需要演示，可以直接启动 GUI，点击生成演示输入，然后加载后端并开始处理，结果会在预览区和输出目录中同步显示。

    ## 第 13 页：测试结果与性能指标（约 1 分 20 秒）
    在测试结果方面，论文中从图像质量、感知质量和处理效率三个维度进行说明。质量指标中，CRA-SRGAN 在 BSD100 上的 PSNR 为 29.45 dB，SSIM 为 0.8623。感知质量指标中，FID 为 3.92，MOS 为 4.25。处理效率方面，4K 图像平均处理时间为 0.32 秒，8K 图像平均处理时间为 0.75 秒。同时，分块处理将 8K 图像显存占用降低到 5.3GB。综合来看，系统不仅完成了功能链路，也具备较好的结果展示能力和工程可用性。

    ## 第 14 页：系统测试与稳定性验证（约 1 分钟）
    系统测试主要分为功能测试和稳定性测试。功能测试覆盖了图像读取、掩码处理、模型加载、图像修复、超分增强、结果保存和图表生成。稳定性测试主要关注异常场景，例如模型路径错误、掩码缺失、图像格式不支持、增强阶段异常等情况。系统会通过日志和弹窗给出提示，而不是直接崩溃。GUI 方面，通过线程方式执行耗时任务，模型加载和图像处理时界面仍能保持响应。批处理方面，个别样本失败不会中断整个批次。

    ## 第 15 页：创新点、不足与后续工作（约 1 分 20 秒）
    本课题的主要完成点有三个。第一，把 CRA 与 SRGAN 从两个独立模型整合成了一条完整的联合图像复原流程。第二，在工程层面加入了分块推理、自动检查点发现、GUI、批处理和掩码分析，使系统具备可运行、可展示、可维护的特点。第三，针对真实运行环境增加了 fallback 机制，提高了演示和部署的稳定性。当然，系统仍存在不足，比如当前采用串联式流程，前一阶段误差可能传递到后一阶段；复杂纹理和极端退化样本的泛化验证还不够充分。后续可以探索更紧密的联合建模，引入轻量化注意力或扩散模型思想，并扩充真实场景测试数据。

    ## 第 16 页：结束页（约 20 秒）
    以上就是我的毕业设计汇报。总体来说，本系统围绕图像缺损和低分辨率这两个实际问题，完成了从模型整合、流程封装、图形界面到测试验证的完整实现。感谢各位老师的聆听，请各位老师批评指正。
    """
).strip()


def build_script():
    md_path = OUT / "答辩演讲稿_约15分钟.md"
    md_path.write_text(SCRIPT + "\n", encoding="utf-8")

    doc = Document()
    styles = doc.styles
    styles["Normal"].font.name = "Microsoft YaHei"
    styles["Normal"].font.size = DocxPt(11)
    for block in SCRIPT.splitlines():
        line = block.strip()
        if not line:
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        else:
            para = doc.add_paragraph(line)
            para.paragraph_format.first_line_indent = DocxPt(22)
            para.paragraph_format.line_spacing = 1.25
    docx_path = OUT / "答辩演讲稿_约15分钟.docx"
    doc.save(docx_path)
    return md_path, docx_path


if __name__ == "__main__":
    deck = build_deck()
    script, docx = build_script()
    print(deck)
    print(script)
    print(docx)
